#!/usr/bin/env python3
"""
聖路加ICUガイドサイト 紹介スライドジェネレーター
Usage: python3 slides/generate_slides.py
Output: slides/output/icu-guide-intro.pptx
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── カラー定義 ────────────────────────────────────────────────
BLUE   = RGBColor(0x1d, 0x4e, 0xd8)   # #1d4ed8
GREEN  = RGBColor(0x2D, 0x7A, 0x4F)   # #2D7A4F
NAVY   = RGBColor(0x1e, 0x3a, 0x5f)   # #1e3a5f
DARK   = RGBColor(0x1e, 0x29, 0x3b)   # #1e293b
MUTED  = RGBColor(0x47, 0x55, 0x69)   # #475569
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT  = RGBColor(0xF1, 0xF5, 0xF9)   # #f1f5f9
BORDER = RGBColor(0xE2, 0xE8, 0xF0)   # #e2e8f0

# ── スライドサイズ (16:9) ─────────────────────────────────────
W = Cm(33.87)
H = Cm(19.05)

# ── パス ─────────────────────────────────────────────────────
ROOT = Path(__file__).parent
SS   = ROOT / "screenshots"
OUT  = ROOT / "output" / "icu-guide-intro.pptx"


# ════════════════════════════════════════════════════════════
# ヘルパー関数
# ════════════════════════════════════════════════════════════

def new_prs() -> Presentation:
    prs = Presentation()
    prs.slide_width  = W
    prs.slide_height = H
    return prs


def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # 6 = Blank


def add_rect(slide, x, y, w, h, color: RGBColor, no_border: bool = True):
    """塗りつぶし矩形を追加。"""
    s = slide.shapes.add_shape(1, x, y, w, h)   # 1 = RECTANGLE
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if no_border:
        s.line.fill.background()
    return s


def add_text(slide, text: str, x, y, w, h,
             size: float, bold: bool = False,
             color: RGBColor = DARK,
             align=PP_ALIGN.LEFT,
             font: str = "Meiryo") -> None:
    """1段落テキストボックスを追加。"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size  = Pt(size)
    run.font.bold  = bold
    run.font.color.rgb = color
    run.font.name  = font


def add_multiline(slide, lines: list, x, y, w, h,
                  size: float, bold: bool = False,
                  color: RGBColor = DARK,
                  align=PP_ALIGN.LEFT,
                  font: str = "Meiryo",
                  line_spacing_pt: float = 4) -> None:
    """
    複数行テキストボックス。
    lines は str のリスト、または (str, bold, color) タプルのリスト。
    """
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, bld, col = item, bold, color
        else:
            txt, bld, col = item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if i > 0 and line_spacing_pt:
            p.space_before = Pt(line_spacing_pt)
        run = p.add_run()
        run.text = txt
        run.font.size  = Pt(size)
        run.font.bold  = bld
        run.font.color.rgb = col
        run.font.name  = font


def add_img(slide, filename: str, x, y, w, h) -> None:
    """
    screenshots/ から画像を挿入。
    ファイルが存在しない場合はプレースホルダー矩形を描く。
    """
    path = SS / filename
    if path.exists():
        slide.shapes.add_picture(str(path), x, y, w, h)
    else:
        add_rect(slide, x, y, w, h, LIGHT)
        add_rect(slide, x, y, w, h, BORDER, no_border=False)
        add_text(slide,
                 f"[スクショ: {filename}]",
                 x + Cm(0.5), y + (h / 2) - Pt(8),
                 w - Cm(1), Cm(1),
                 size=9, color=MUTED, align=PP_ALIGN.CENTER)


def page_num(slide, n: int, total: int = 15) -> None:
    """右下にページ番号。"""
    add_text(slide, f"{n} / {total}",
             W - Cm(3), H - Cm(1.2),
             Cm(2.5), Cm(0.8),
             size=9, color=BORDER, align=PP_ALIGN.RIGHT)


# ════════════════════════════════════════════════════════════
# レイアウト① 表紙
# ════════════════════════════════════════════════════════════

def slide_cover(prs: Presentation) -> None:
    """スライド1 — 表紙"""
    sl = blank_slide(prs)

    # 背景：ネイビー全面
    add_rect(sl, 0, 0, W, H, NAVY)

    # 左側ブルー帯
    add_rect(sl, 0, 0, W * 0.6, H, BLUE)

    # 右側グリーン帯
    add_rect(sl, W * 0.4, 0, W * 0.6, H, GREEN)

    # 濃紺オーバーレイ（深みを出す）
    add_rect(sl, 0, 0, W, H, NAVY)

    # 左アクセント：ブルー縦バー
    add_rect(sl, 0, 0, Cm(0.8), H, BLUE)

    # 下部グリーンバー
    add_rect(sl, 0, H - Cm(1.8), W, Cm(1.8), GREEN)

    # 施設名（小）
    add_text(sl, "聖路加国際病院 ICU",
             Cm(3), Cm(5.5), Cm(22), Cm(1.2),
             size=14, bold=False, color=RGBColor(0xA0, 0xB8, 0xD8))

    # メインタイトル
    add_text(sl, "ICU ローテーターガイド",
             Cm(3), Cm(6.8), Cm(22), Cm(3),
             size=36, bold=True, color=WHITE)

    # サブタイトル
    add_text(sl, "— サイト紹介 —",
             Cm(3), Cm(9.8), Cm(22), Cm(1.5),
             size=18, bold=False, color=RGBColor(0xA0, 0xB8, 0xD8))

    # URL
    add_text(sl, "https://sekisato59.github.io/st-lukes-icu/",
             Cm(3), Cm(12), Cm(22), Cm(1),
             size=11, bold=False, color=RGBColor(0x7A, 0xA8, 0xD8))

    # QRプレースホルダー（右下）
    add_img(sl, "qr_icu_guide.png",
            W - Cm(6.5), H - Cm(7.5), Cm(5.5), Cm(5.5))


# ════════════════════════════════════════════════════════════
# レイアウト② セクション扉
# ════════════════════════════════════════════════════════════

def slide_section(prs: Presentation, n: int,
                  part_label: str, title: str,
                  body_lines: list,
                  screenshot: str = "") -> None:
    """
    スライド2・3・4・12 — セクション扉
    body_lines: str のリスト（箇条書き）
    screenshot: screenshots/ 内のファイル名（空白可）
    """
    sl = blank_slide(prs)

    # 白背景
    add_rect(sl, 0, 0, W, H, WHITE)

    # 左縦カラーバー（ブルー上半分・グリーン下半分）
    add_rect(sl, 0, 0, Cm(0.5), H * 0.5, BLUE)
    add_rect(sl, 0, H * 0.5, Cm(0.5), H * 0.5, GREEN)

    # PARTラベル
    add_text(sl, part_label,
             Cm(1.5), Cm(3), Cm(15), Cm(0.9),
             size=11, bold=True, color=BLUE)

    # タイトル
    add_text(sl, title,
             Cm(1.5), Cm(4), Cm(18), Cm(3.5),
             size=28, bold=True, color=DARK)

    # 本文（箇条書き）
    bullet_lines = [f"  ▸  {line}" for line in body_lines]
    add_multiline(sl, bullet_lines,
                  Cm(1.5), Cm(8), Cm(18), Cm(7),
                  size=14, color=MUTED, line_spacing_pt=6)

    # 右側スクショ
    if screenshot:
        add_img(sl, screenshot, Cm(20), Cm(2.5), Cm(12.5), Cm(13))
    else:
        add_rect(sl, Cm(20), Cm(3), Cm(12), Cm(12), LIGHT)
        add_text(sl, "ICU",
                 Cm(20), Cm(7.5), Cm(12), Cm(3),
                 size=40, align=PP_ALIGN.CENTER, color=GREEN)

    page_num(sl, n)


# ════════════════════════════════════════════════════════════
# レイアウト③ スクショ中心
# ════════════════════════════════════════════════════════════

def slide_feature(prs: Presentation, n: int,
                  part_label: str, title: str,
                  bullets: list,
                  ss_left: str, ss_right: str) -> None:
    """
    スライド5〜11 — 機能紹介（スクショ2枚）
    bullets: str のリスト（3〜4項目）
    ss_left / ss_right: screenshots/ 内のファイル名
    """
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, W, H, WHITE)

    # ── ヘッダー帯 ─────────────────────────────────────────
    HEADER_H = Cm(1.8)
    add_rect(sl, 0, 0, W * 0.5, HEADER_H, BLUE)
    add_rect(sl, W * 0.5, 0, W * 0.5, HEADER_H, GREEN)

    add_text(sl, part_label,
             Cm(0.8), Cm(0.2), Cm(6), HEADER_H,
             size=10, bold=True, color=RGBColor(0xA0, 0xC8, 0xFF))

    add_text(sl, title,
             Cm(5), Cm(0.2), Cm(22), HEADER_H,
             size=15, bold=True, color=WHITE)

    # ── 左1/3: 説明テキスト ────────────────────────────────
    TEXT_X = Cm(0.8)
    TEXT_Y = HEADER_H + Cm(0.8)
    TEXT_W = Cm(9.5)

    bullet_lines = [f"▸  {b}" for b in bullets]
    add_multiline(sl, bullet_lines,
                  TEXT_X, TEXT_Y, TEXT_W, Cm(12),
                  size=13, color=MUTED, line_spacing_pt=8)

    # 誘導ラベル
    add_rect(sl, TEXT_X, H - Cm(2.5), TEXT_W, Cm(1.2),
             RGBColor(0xDB, 0xEA, 0xFE))
    add_text(sl, "→ 詳しくはサイトで",
             TEXT_X + Cm(0.3), H - Cm(2.5), TEXT_W - Cm(0.5), Cm(1.2),
             size=10, bold=True, color=BLUE)

    # ── 右2/3: スクショ2枚 ────────────────────────────────
    SS_X  = Cm(10.8)
    SS_Y  = HEADER_H + Cm(0.6)
    SS_W  = Cm(10.8)
    SS_H  = Cm(14.5)

    add_img(sl, ss_left,  SS_X,             SS_Y, SS_W, SS_H)
    add_img(sl, ss_right, SS_X + SS_W + Cm(0.5), SS_Y, SS_W, SS_H)

    page_num(sl, n)


# ════════════════════════════════════════════════════════════
# レイアウト④ 締め
# ════════════════════════════════════════════════════════════

def slide_closing(prs: Presentation, n: int,
                  label: str, title: str,
                  body_lines: list,
                  show_qr: bool = False) -> None:
    """スライド13〜15 — 締め"""
    sl = blank_slide(prs)
    add_rect(sl, 0, 0, W, H, WHITE)

    # 左縦カラーバー（グリーン上・ブルー下）
    add_rect(sl, 0, 0, Cm(0.5), H * 0.5, GREEN)
    add_rect(sl, 0, H * 0.5, Cm(0.5), H * 0.5, BLUE)

    # ラベル
    add_text(sl, label,
             Cm(1.5), Cm(3), Cm(20), Cm(0.9),
             size=11, bold=True, color=GREEN)

    # タイトル
    add_text(sl, title,
             Cm(1.5), Cm(4), Cm(20), Cm(2.5),
             size=26, bold=True, color=DARK)

    # 本文
    add_multiline(sl, body_lines,
                  Cm(1.5), Cm(7), Cm(20), Cm(9),
                  size=14, color=MUTED, line_spacing_pt=8)

    # QR（最終スライドのみ）
    if show_qr:
        add_img(sl, "qr_icu_guide.png",
                W - Cm(7.5), Cm(4), Cm(6.5), Cm(6.5))
        add_text(sl,
                 "https://sekisato59.github.io/st-lukes-icu/",
                 W - Cm(8), Cm(11), Cm(7.5), Cm(1),
                 size=9, color=BLUE)

    page_num(sl, n)


# ════════════════════════════════════════════════════════════
# 全スライド定義
# ════════════════════════════════════════════════════════════

def build(prs: Presentation) -> None:

    # ── スライド 1: 表紙 ────────────────────────────────────
    slide_cover(prs)

    # ── スライド 2: 困りごと ───────────────────────────────
    slide_section(prs, 2,
        part_label="PART 0 — はじめに",
        title="ICUローテーションの\n「困りごと」",
        body_lines=[
            "情報が散在していて、どこを見ればいいかわからない",
            "論文・ガイドラインの調べ方が手探り",
            "計算ツールやスコアをいちいち調べている",
            "こんな悩みを解決するために作りました",
        ],
        screenshot="01_top-hero.png")

    # ── スライド 3: 解決できること ─────────────────────────
    slide_section(prs, 3,
        part_label="PART 0 — はじめに",
        title="このサイトで\n解決できること",
        body_lines=[
            "① 病棟情報 — スケジュール・ルール・スタッフ",
            "② 学習コンテンツ — 論文GL・疾患マニュアル・動画",
            "③ 便利ツール — 計算機・スコア・カルテ略語",
            "④ 全文検索 — キーワード1つで即アクセス",
        ],
        screenshot="")

    # ── スライド 4: サイト全体像 ───────────────────────────
    slide_section(prs, 4,
        part_label="PART 0 — はじめに",
        title="サイト全体像",
        body_lines=[
            "トップページから全コンテンツにアクセス可能",
            "PCデスクトップ最適化（スマホも利用可）",
            "随時更新中（論文GL・疾患マニュアルを毎月追加）",
        ],
        screenshot="02_top-overview.png")

    # ── スライド 5: 当院ICUについて ─────────────────────────
    slide_feature(prs, 5,
        part_label="PART 1 — 病棟情報",
        title="当院ICUについて",
        bullets=[
            "ICU概要・スタッフ紹介",
            "1日・週間のスケジュール",
            "ICU独自のルール・診療方針",
            "ローテ開始前のTo Doリスト",
        ],
        ss_left="03_about-icu.png",
        ss_right="04_pre-rotation.png")

    # ── スライド 6: 注目論文・ガイドラインまとめ ─────────────
    slide_feature(prs, 6,
        part_label="PART 2 — 学習コンテンツ",
        title="注目論文・ガイドラインまとめ",
        bullets=[
            "ICU関連の最新GL・注目論文を臓器系統別に整理",
            "原著PDFに基づく要約（ハルシネーションなし）",
            "stat-gridで統計データを視覚化",
            "循環器・呼吸器・感染症・腎など10系統以上収録",
        ],
        ss_left="05_articles-gl.png",
        ss_right="06_articles-gl-detail.png")

    # ── スライド 7: 疾患マニュアル ─────────────────────────
    slide_feature(prs, 7,
        part_label="PART 2 — 学習コンテンツ",
        title="疾患マニュアル",
        bullets=[
            "収録済み論文・GLの知見を疾患ごとに統合",
            "後期研修医・フェロー向け総説ページ集",
            "敗血症・ARF・AKI・心不全など多数収録",
            "GLページへのクロスリンクで深掘り可能",
        ],
        ss_left="07_disease-topics.png",
        ss_right="08_disease-detail.png")

    # ── スライド 8: ICU PASSPORT・動画講座 ─────────────────
    slide_feature(prs, 8,
        part_label="PART 2 — 学習コンテンツ",
        title="ICU PASSPORT・動画講座",
        bullets=[
            "4種のミッションで学習モチベーションをキープ",
            "ICU動画講座集 — 50本以上を自分のペースで視聴",
            "ID×ICU Conference — 毎週の感染症レクチャー録画",
            "学びたい内容を、学びたい時に、学びたいだけ",
        ],
        ss_left="09_icu-passport.png",
        ss_right="10_video-lectures.png")

    # ── スライド 9: 計算ツール・スコア ─────────────────────
    slide_feature(prs, 9,
        part_label="PART 3 — 便利ツール",
        title="計算ツール・スコア一覧",
        bullets=[
            "抗菌薬投与量 — CrCl・体重を入れるだけ",
            "栄養・輸血量 → 上昇予測計算",
            "APACHE II / SOFA-2 / RASS / CAM-ICU",
            "LDL目標値チェッカー（JAS 2022準拠）",
        ],
        ss_left="11_tools-calc.png",
        ss_right="12_tools-score.png")

    # ── スライド 10: カルテ略語・細菌マップ ─────────────────
    slide_feature(prs, 10,
        part_label="PART 3 — 便利ツール",
        title="カルテ略語チェッカー・細菌マップ",
        bullets=[
            "略語を自動スペルアウト・禁止表記チェック",
            "オフライン動作（院内ネットワークでも使用可）",
            "細菌マップ — グラム染色×形態×嫌気/好気で整理",
            "代表感染症・第一選択薬・関連GLへリンク",
        ],
        ss_left="13_karte-abbr.png",
        ss_right="14_bacteria-map.png")

    # ── スライド 11: 検索機能① ─────────────────────────────
    slide_feature(prs, 11,
        part_label="PART 4 — 検索機能",
        title="検索機能の使い方 ①",
        bullets=[
            "ナビバー右端の「検索」をクリック",
            "キーワードを入力（例：抗菌薬、透析、ルール）",
            "タイトル・キーワード・説明文を横断検索",
            "結果はリアルタイムに絞り込まれる",
        ],
        ss_left="15_search-box.png",
        ss_right="16_search-results.png")

    # ── スライド 12: 検索機能② ─────────────────────────────
    slide_section(prs, 12,
        part_label="PART 4 — 検索機能",
        title="検索機能の使い方 ②\n— 該当ページへジャンプ",
        body_lines=[
            "検索結果をクリックすると対象ページに直接ジャンプ",
            "各ページはナビバーの統一メニューから行き来可能",
            "「コンテンツをすべて検索」から全ページを一覧できる",
        ],
        screenshot="17_search-jump.png")

    # ── スライド 13: 対象読者・更新 ────────────────────────
    slide_closing(prs, 13,
        label="まとめ",
        title="こんな先生におすすめ",
        body_lines=[
            "▸  聖路加ICUにローテーション中の初期・後期研修医",
            "▸  ICU集中治療を学びたいすべての研修医・フェロー",
            "▸  最新の論文・GLを効率よくキャッチアップしたい方",
            "",
            "更新について",
            "▸  論文GL・疾患マニュアルを随時追加中",
            "▸  トップページ「新着コンテンツ」で最新情報を確認",
        ])

    # ── スライド 14: さいごに ───────────────────────────────
    slide_closing(prs, 14,
        label="さいごに",
        title="制作の思い",
        body_lines=[
            "ICUローテーションは情報量が多く、戸惑う先生が多い。",
            "",
            "「必要な情報に、すぐアクセスできる環境」を整えることで",
            "研修医が学習・診療に集中できるよう願って制作しました。",
            "",
            "フィードバック・追加リクエストはいつでも歓迎です。",
        ])

    # ── スライド 15: アクセス方法 ─────────────────────────
    slide_closing(prs, 15,
        label="ACCESS",
        title="アクセス方法",
        body_lines=[
            "URL:",
            "https://sekisato59.github.io/st-lukes-icu/",
            "",
            "▸  Chrome / Safari / Edge でそのまま開けます",
            "▸  ブックマーク・ホーム画面追加を推奨",
            "",
            "制作: 聖路加国際病院 ICU 関谷",
        ],
        show_qr=True)


# ════════════════════════════════════════════════════════════
# エントリーポイント
# ════════════════════════════════════════════════════════════

if __name__ == "__main__":
    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs = new_prs()
    build(prs)
    prs.save(str(OUT))
    print(f"✅  生成完了: {OUT}")
    print(f"   スライド数: {len(prs.slides)}")
